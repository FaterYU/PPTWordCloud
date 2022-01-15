from pptx import Presentation
from wordcloud import WordCloud


def Read_ppt(path_to_presentation):
    prs = Presentation(path_to_presentation)
    text_runs = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs += run.text
    return text_runs


def Creat_wc(text):
    font = r'C:\Windows\Fonts\simfang.ttf'
    wc = WordCloud(collocations=False, font_path=font, width=1400,
                   height=1400, margin=2).generate(text.lower())
    return wc
